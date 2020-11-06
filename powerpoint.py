import win32com.client
import logging
from pptx import Presentation
from pptx.dml.color import RGBColor
from PIL import ImageColor
import os

PPT_PATH = r'D:\Project\commerce\parser_powerpoint\init_pptx\Fashion.pptx'
PPT_CHANGED_PATH = r'D:\Project\commerce\parser_powerpoint\init_pptx\new.pptx'

logger = logging.getLogger("parser_loger")
logger.setLevel(logging.INFO)
console_headler = logging.StreamHandler()
formatter = logging.Formatter('[%(levelname)s] %(asctime)s %(lineno)d : %(message)s')
console_headler.setFormatter(formatter)
logger.addHandler(console_headler)


# Types 17 - Text box
# Types 14 - Placeholder
# Types 16 - Media


class Slide:
    SPEED = {
        1: 'Медленно',
        2: 'Средне',
        3: 'Быстро',
    }

    def __init__(self, ppt, slide_com_object):
        self.ppt = ppt
        self.slide_com_object = slide_com_object
        self.shapes = [Shape(self, shape) for shape in self.slide_com_object.Shapes]
        self.id = self.slide_com_object.SlideIndex
        self.texts = self.get_texts()
        self.images = self.get_pictures()
        self.videos = self.get_videos()
        self.audios = self.get_audios()
        self.speed = self.get_speed()

    def get_texts(self) -> list:
        logger.info('Получение текстовой информации слайда №{}'.format(self.id))
        text_frames = []
        for shape in self.shapes:
            if shape.Type == 17:
                text_frame = shape.TextFrame
                if text_frame.HasText:
                    text_frames.append(FrameText(text_frame))
        return text_frames

    def get_pictures(self) -> list:
        logger.info('Получение картинок слайда №{}'.format(self.id))
        images = []
        for shape in self.shapes:
            if shape.Type == 14:
                images.append(ShapeImage(self, shape))
        return images

    def get_videos(self):
        logger.info('Получение видео слайда №{}'.format(self.id))
        videos = []
        for shape in self.shapes:
            if shape.Type == 16:
                if shape.MediaType == 3:
                    videos.append(ShapeVideo(self, shape))
        return videos

    def get_audios(self):
        logger.info('Получение аудио слайда №{}'.format(self.id))
        audios = []
        for shape in self.shapes:
            if shape.Type == 16:
                if shape.MediaType == 2:
                    audios.append(ShapeAudio(self, shape))
        return audios

    def change_speed(self, speed):
        logger.info('Изменение скоррости слайда №{} с {} на {} '.format(self.id,
                                                                        self.speed,
                                                                        self.SPEED[speed]))
        self.slide_com_object.SlideShowTransition.Speed = speed
        print(self.slide_com_object.SlideShowTransition.Speed)

    def get_speed(self):
        logger.info('Получение продолжительности слайда №{}'.format(self.id))
        return self.SPEED[self.slide_com_object.SlideShowTransition.Speed]

    def change_background_color(self, rgb_touple):
        file = os.path.dirname(__file__) + '\_.pptx'
        self.ppt.save_as(file)
        self.ppt.close()
        ppt_with_changing_colors = PPTColors(file)
        ppt_with_changing_colors.slides[self.id - 1].change_background_color(rgb_touple)
        ppt_with_changing_colors.save_as(file)
        self.ppt.__init__(file)


class Shape:
    def __init__(self, slide, shape_object):
        self.slide = slide
        self.shape_object = shape_object
        self.Type = shape_object.Type
        self.MediaType = None
        if self.Type == 16:
            self.MediaType = shape_object.MediaType
        self.TextFrame = None
        if self.Type == 17:
            self.TextFrame = shape_object.TextFrame
        self.Left = shape_object.Left
        self.Top = shape_object.Top
        self.Width = shape_object.Width
        self.Height = shape_object.Height
        self.Delete = shape_object.Delete

    def change_color(self, rgb_tuple):
        file = os.path.dirname(__file__) + '\_.pptx'
        self.slide.ppt.save_as(file)
        self.slide.ppt.close()
        ppt_with_changing_colors = PPTColors(file)
        ppt_with_changing_colors.slides[self.slide.id - 1].shapes[self.slide.shapes.index(self)].change_color(rgb_tuple)
        ppt_with_changing_colors.save_as(file)
        self.slide.ppt.__init__(file)


class FrameText:
    def __init__(self, text_frame_com_object):
        self.text_frame_com_object = text_frame_com_object

    def __repr__(self):
        return self.text_frame_com_object.TextRange.Text

    def change_text(self, text):
        logger.info('Изменение теста {} на {}'.format(self, text))
        self.text_frame_com_object.TextRange.Text = text

    def show_text(self):
        print(self)


class ShapeMedia:
    def __init__(self, slide, shape):
        self.slide = slide
        self.shape = shape

    def change_media(self, path, left=None, top=None, width=None, height=None):
        if not left:
            left = self.shape.Left
        if not top:
            top = self.shape.Top
        if not width:
            width = self.shape.Width
        if not height:
            height = self.shape.Height
        self.shape.Delete()
        print(self.slide.audios)
        self.slide.slide_com_object.Shapes.AddMediaObject(FileName=path,
                                                          Left=left,
                                                          Top=top,
                                                          Width=width,
                                                          Height=height)


class ShapeImage:
    def __init__(self, slide, shape):
        self.slide = slide
        self.shape = shape

    def change_image(self, image_path):
        logger.info('Изменение изображения в {} слайде'.format(self.slide.id))
        left = self.shape.Left
        top = self.shape.Top
        width = self.shape.Width
        height = self.shape.Height
        self.shape.Delete()
        self.slide.slide_com_object.Shapes.AddPicture(FileName=image_path,
                                                      LinkToFile=False,
                                                      SaveWithDocument=True,
                                                      Left=left,
                                                      Top=top,
                                                      Width=width,
                                                      Height=height)


class ShapeVideo(ShapeMedia):
    def __init__(self, slide, shape):
        super().__init__(slide, shape)

    def change_video(self, video_path, left=None, top=None, width=None, height=None):
        logger.info('Изменение видео в {} слайде'.format(self.slide.id))
        self.change_media(video_path, left, top, width, height)


class ShapeAudio(ShapeMedia):
    def __init__(self, slide, shape):
        super().__init__(slide, shape)

    def change_audio(self, audio_path, left=None, top=None, width=None, height=None):
        logger.info('Изменение аудио в {} слайде'.format(self.slide.id))
        self.change_media(audio_path, left, top, width, height)


class PPT:
    def __init__(self, ppt_path):
        logger.info('Открытие презентации {}'.format(ppt_path))
        self.app = win32com.client.Dispatch("Powerpoint.Application")
        self.ppt_path = ppt_path
        self.ppt_com_object = self.app.Presentations.Open(self.ppt_path)
        self.slides = self.get_slides()
        self.ppt_pptx_python_object = PPTColors(ppt_path)

    def get_slides(self) -> list:
        logger.info('Получение слайдов')
        return [Slide(self, slide_com_object) for slide_com_object in self.ppt_com_object.Slides]

    def close(self):
        self.ppt_com_object.Close()

    def quit(self):
        self.app.Quit()

    def duplicate_slide(self, index_slide, index_copy_place):
        logger.info('Копирование слайда №{} на место слайда №{}'.format(index_slide, index_copy_place))
        self.ppt_com_object.Slides[index_slide - 1].Copy()
        self.ppt_com_object.Slides.Paste(index_copy_place)

    def save_as(self, file_name):
        logger.info('Сохранение презентации {}'.format(file_name))
        self.ppt_com_object.SaveAs(file_name)

    def show_info(self):
        info_main = 'Слайд #{}:'
        info_shade = '\tФормы : {}'
        info_shade_color = '\tЦвета форм : {}'
        info_background = '\tФон : {}'
        info_text_blocks = '\tТекстовые блоки : {}'
        info_pictures = '\tКартинки : {}'
        info_videos = '\tВидео : {}'
        info_audios = '\tАудио : {}'
        info_speed_slide = '\tСкорость слайда: {}'
        for slide_index in range(len(self.slides)):
            print(info_main.format(slide_index + 1))
            print(info_shade.format([shape.type for shape in self.ppt_pptx_python_object.slides[slide_index].shapes]))
            print(info_shade_color.format(
                [shape.color for shape in self.ppt_pptx_python_object.slides[slide_index].shapes]))
            print(info_background.format(
                'тип - {}, цвет - {}'.format(self.ppt_pptx_python_object.slides[slide_index].background,
                                             self.ppt_pptx_python_object.slides[slide_index].background_color)))
            print(info_text_blocks.format([str(text)[:4] + '...' for text in self.slides[slide_index].texts]))
            print(info_pictures.format(
                ['img' + str(image_index) for image_index in range(len(self.slides[slide_index].images))]))
            print(info_videos.format(
                ['video' + str(video_index) for video_index in range(len(self.slides[slide_index].videos))]))
            print(info_audios.format(
                ['audio' + str(audio_index) for audio_index in range(len(self.slides[slide_index].audios))]))
            print(info_speed_slide.format(self.slides[slide_index].speed))


class SlideColors:
    FILL_TYPES = {
        5: 'Прозрачный',
        3: 'Градиент',
        -2: 'Группа',
        2: 'Шаблонн',
        6: 'Картинка',
        1: 'Заливка',
        4: 'Текстура',
    }

    def __init__(self, slide_object):
        self.slide_object = slide_object
        self.id = self.slide_object.slide_id
        self.shapes = [ShapeColors(shape) for shape in self.slide_object.shapes]
        self.background = self.FILL_TYPES[self.slide_object.background.fill.type]
        self.background_color = None
        self.check_background_color()

    def check_background_color(self):
        if self.slide_object.background.fill.type == 1:
            if self.slide_object.background.fill.fore_color.type == 1:
                self.background_color = ImageColor.getcolor("#" + str(self.slide_object.background.fill.fore_color.rgb),
                                                            "RGB")

    def change_background_color(self, rgb_tuple):
        logger.info('Изменение цвета фона {}  на цвет {}'.format(self.background_color, rgb_tuple))
        self.slide_object.background.fill.solid()
        self.slide_object.background.fill.fore_color.rgb = RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])


class ShapeColors:
    SHAPES_TYPES = {
        1: 'Авто форма',
        5: 'Свободная форма',
        6: 'Группа',
        17: 'Текст бокс',
        13: 'Картинка',
        14: 'Заполнитель',
        16: 'Медиа',
    }

    def __init__(self, shape_object):
        self.shape_object = shape_object
        self.type = self.SHAPES_TYPES[self.shape_object.shape_type]
        self.color = None
        self.check_color()

    def check_color(self):
        if self.shape_object.shape_type in (6, 13, 14, 16):
            return
        if self.shape_object.fill.type in (3, 5):
            return
        if self.shape_object.fill.fore_color.type == 1:
            self.color = ImageColor.getcolor("#" + str(self.shape_object.fill.fore_color.rgb),
                                             "RGB")

    def change_color(self, rgb_tuple):
        if self.shape_object.shape_type in (6, 13, 14, 16):
            logger.info('Изменение цвета у данных типов форм невозможно')
            return
        logger.info('Изменение цвета формы {}  на цвет {}'.format(self.color, rgb_tuple))
        self.shape_object.fill.solid()
        self.shape_object.fill.fore_color.rgb = RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])


class PPTColors:
    def __init__(self, ppt_path):
        logger.info('Открытие презентации {}'.format(ppt_path))
        self.ppt = Presentation(pptx=ppt_path)
        self.slides = [SlideColors(slide) for slide in self.ppt.slides]

    def save_as(self, file_name):
        logger.info('Сохранение презентации {}'.format(file_name))
        self.ppt.save(file_name)
